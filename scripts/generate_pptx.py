#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script to generate a PPTX product showcase from images in the repository.
It groups images by top-level folders and creates slides with a cute, rounded style.
Run: python scripts/generate_pptx.py
This script is designed to run inside the repository (CI or locally).
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
IMAGE_DIRS = [p for p in ROOT.iterdir() if p.is_dir() and not p.name.startswith('.') and p.name != '.github' and p.name != 'scripts' and p.name != 'PPT']
OUT_DIR = ROOT / 'PPT'
OUT_DIR.mkdir(exist_ok=True)
OUT_FILE = OUT_DIR / '鲍勃熊产品展示.pptx'

# Presentation defaults
prs = Presentation()
SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(14), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = '鲍勃熊产品展示'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 178)

    p2 = tf.add_paragraph()
    p2.text = 'Bobbear Product Showcase'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(120, 120, 120)

    tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(0.6))
    tf2 = tx2.text_frame
    p = tf2.paragraphs[0]
    p.text = f'作者: {os.getenv("GITHUB_ACTOR", "lrbgo")}'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)

def add_section_slide(prs, title_zh, title_en):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(14), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title_zh
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 153, 102)

    p2 = tf.add_paragraph()
    p2.text = title_en
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(120, 120, 120)

def add_image_slide(prs, img_path, caption_zh, caption_en):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    margin = Inches(0.8)
    max_w = SLIDE_WIDTH - margin * 2
    max_h = SLIDE_HEIGHT - margin * 2 - Inches(1.2)
    with Image.open(img_path) as im:
        width_px, height_px = im.size
        dpi = im.info.get('dpi', (96,96))[0]
        img_w_in = width_px / dpi if dpi else width_px / 96
        img_h_in = height_px / dpi if dpi else height_px / 96

    if img_w_in == 0 or img_h_in == 0:
        img_w_in = width_px / 96
        img_h_in = height_px / 96

    max_w_in = SLIDE_WIDTH.inches - margin.inches * 2
    max_h_in = SLIDE_HEIGHT.inches - margin.inches * 2 - 1.2
    scale = min(max_w_in / img_w_in, max_h_in / img_h_in, 1.0)
    disp_w = img_w_in * scale
    disp_h = img_h_in * scale

    left = (SLIDE_WIDTH.inches - disp_w) / 2
    top = (SLIDE_HEIGHT.inches - disp_h) / 2 - 0.4

    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(disp_w), height=Inches(disp_h))

    tx = slide.shapes.add_textbox(Inches(1.2), SLIDE_HEIGHT - Inches(1.1), SLIDE_WIDTH - Inches(2.4), Inches(0.9))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = caption_zh
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(60, 60, 60)

    p2 = tf.add_paragraph()
    p2.text = caption_en
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(120, 120, 120)

def safe_name_from_path(p: Path):
    return p.stem

add_cover(prs)

for folder in sorted(IMAGE_DIRS):
    zh_section = folder.name
    en_section = ' '.join([part.capitalize() for part in folder.name.split()])
    add_section_slide(prs, zh_section, en_section)
    imgs = [f for f in sorted(folder.iterdir()) if f.is_file() and f.suffix.lower() in ('.png', '.jpg', '.jpeg')]
    for img in imgs:
        name = safe_name_from_path(img)
        caption_zh = name
        caption_en = f'EN: {name}'
        try:
            add_image_slide(prs, img, caption_zh, caption_en)
        except Exception as e:
            print(f'Failed to add image {img}: {e}')

slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(14), Inches(2))
tf = tx.text_frame
p = tf.paragraphs[0]
p.text = '谢谢观看'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 102, 178)

p2 = tf.add_paragraph()
p2.text = 'Thank you'
p2.font.size = Pt(18)

prs.save(str(OUT_FILE))
print('Saved PPTX to', OUT_FILE)
